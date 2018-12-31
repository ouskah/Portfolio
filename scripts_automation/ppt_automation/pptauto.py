from pptx import Presentation
import pandas as pd
import names
import random


# Création d'un set de connées fictives qui imitent les fichier csv de la mission

# Création d'une liste de 100 noms aléatoires
liste_noms = [names.get_full_name() for x in range(0,100)]
# Création d'une liste de 100 notes aléatoires (allant de 0 à 10)
liste_notes = [random.randint(0,10) for x in range(0,100)]
# Liste finale qui contiendra l'ensemble des commentaires des utilisateurs
liste_comment = []
# Intégration d'un premier pattern de commentaires qui correspondent au profil utilisateur de "promoteur"
pattern_comments_promoteurs = ['Excellent !', 'Parfait !', 'Génial !', 'j\'aime beaucoup',
                               'c\'est super !', 'Wow', 'Incroyable', 'Fantastique !',
                               'Impressionnant', 'Magnifique', 'Prodigieux']
# Intégration d'un second pattern de commentaires qui correspondent au profil utilisateur de "passif"
pattern_comments_passifs = ['moyen..', 'pas mal', 'c\'est correct', 'ok',
                            'passable', 'pourquoi pas...admettons..', 'mouais...',
                            'pas ouf comme produit, mais ça passe']
# Intégration d'un troisième pattern de commentaires qui correspondent au profil utilisateur de "detracteur"
pattern_comments_detracteurs = ['nul !', 'je n\'aime pas du tout', 'produit défectueux', 'vraiment pas terrible...',
                                'Attention danger !', 'Ne surtout pas acheter..', 'A ne jamais refaire']

# Ajout d'un commentaire de manière aléatoire en fonction des catégories de notation
# - Un promoteur est un utilisateur qui a attribué une note supérieur à 7 sur 10
# - Un profil passif est un utilisateur qui a attribué une note qui est égale à 5 sur 10
# - Un détracteur est un utilisateur qui a attribué une note qui est supérieure ou égale à 0 et inférieure à 5 sur 10
for y in liste_notes:
    if y > 7:
        liste_comment.append(random.choice(pattern_comments_promoteurs))
    elif y == 5:
        liste_comment.append(random.choice(pattern_comments_passifs))
    elif (y > 0) and (y < 5):
        liste_comment.append(random.choice(pattern_comments_detracteurs))
    elif y == 0:
        liste_comment.append(random.choice(pattern_comments_detracteurs))
    else:
        liste_comment.append(" ")

# Vérification de l'uniformisation des fichier (same length)
# print(len(liste_comment))
# print(len(liste_notes))
# print(len(liste_noms))

# Création du dataframe & export en csv ou xls
df = pd.DataFrame({"nom/prenom": liste_noms, "note(0-10)": liste_notes, "liste_comment": liste_comment})
# df.to_csv("avis_utilisateurs_fictifs.csv", index=False)

# Création des trois segments définis par l'entreprise
segment_promoteurs = df.loc[df["note(0-10)"] > 7]
segment_passifs = df.loc[(df["note(0-10)"] > 4) & (df["note(0-10)"] < 8)]
segment_detracteurs = df.loc[(df["note(0-10)"] >= 0) & (df["note(0-10)"] < 5)]
# Taille des différents segments
seg1 = len(segment_promoteurs)
seg2 = len(segment_passifs)
seg3 = len(segment_detracteurs)
# On affiche la taille de chaque segment
# print(seg1)
# print(seg2)
# print(seg3)

# présenter les commentaires des promoteurs (max 12)
def selection_comment():
    liste_comments = {"promoteurs": [], "passifs": [], "detracteurs": []}
    ctps, ctd = 1, 1
    # Choisir les 12 commentaires de promoteurs que l'on souhaite via la liste de commentaires
    while len(liste_comments["promoteurs"]) < 12:
        for x, y in zip(segment_promoteurs.liste_comment, range(1, seg1)):
            # Option : print l'ensemble de la liste des comments avant ? (= premier apercu)
            if x == " ":
                pass
            else:
                print("Commentaire de promoteur n°{} : {}".format(y, x))
                valider = input("Valider ce commentaire ? (O/n) > ")
                if valider.lower() == "o":
                    liste_comments["promoteurs"].append(x)
                if len(liste_comments["promoteurs"]) == 12:
                    break
    # Choisir le commentaire d'utilisateur passif à partir de la liste de commentaires des utilisateurs passifs            
    for x, y in zip(segment_passifs.liste_comment, range(1, seg2)):
        if x == " ":
            pass
        else:
            print("Commentaire d\'utilisateur passif n°{} : {}".format(y, x))
        valider = input("Valider ce commentaire ? (O/n) > ")
        if valider.lower() == "o":
            liste_comments["passifs"].append(x)
            ctps = ctps - 1
            if ctps == 0:
                break
    # Choisir le commentaire de detracteur à partir de la liste de commentaires des detracteurs            
    for x, y in zip(segment_detracteurs.liste_comment, range(1, seg3)):
        if x == " ":
            pass
        else:
            print("Commentaire de detracteur n°{} : {}".format(y, x))
        valider = input("Valider ce commentaire ? (O/n) > ")
        if valider.lower() == "o":
            liste_comments["detracteurs"].append(x)
            ctd = ctd - 1
            if ctd == 0:
                break
                
    return liste_comments["promoteurs"], liste_comments["passifs"], liste_comments["detracteurs"]


liste_seg = [seg1, seg2, seg3]
liste_labels = ["% de promoteurs", "% de passifs", "% de détracteurs"]
liste_comment_prom, liste_comment_pass, liste_comment_detr = selection_comment()

######################## Création du graphique de type donut chart ###############################
import matplotlib.pyplot as plt
import matplotlib as mpl

mpl.rcParams['font.size'] = 9.0
fig = plt.figure(figsize=(15,15))
plt.title('Pourcentage d\'utilisateurs par segment', fontsize=35)
# Pie chart
labels = liste_labels
sizes = [x for x in liste_seg]
#colors
colors = ['#ff9999','#66b3ff','#99ff99']
explode = (0.05,0.05,0.05)
patches, texts, autotexts = plt.pie(sizes, colors = colors, labels=labels, autopct='%1.1f%%', 
                                    startangle=90, pctdistance=0.85, explode = explode)

# Changer dynamiquement la taille de la police
texts[0].set_fontsize(26)
texts[1].set_fontsize(26)
texts[2].set_fontsize(26)
autotexts[0].set_fontsize(26)
autotexts[1].set_fontsize(26)
autotexts[2].set_fontsize(26)

centre_circle = plt.Circle((0,0),0.70,fc='white')
fig = plt.gcf()
fig.gca().add_artist(centre_circle)

plt.tight_layout()
plt.savefig('data_graphiques/donut_chart.png', bbox_inches="tight")

######################################## Création du ppt #########################################
# format demandé :
# slide 1 = - 3 commentaires de promoteurs 
# + 1 commentaire d'un utilisateur passif 
# + 1 commentaire d'un détracteur
# slide intermédiaire = graphique représentant les proportions des commentaires par segment
# slide 3 = 
# - 9 commentaires de promoteur sur une autre diapositive.
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_UNDERLINE

prs = Presentation()

#---------------------------------    Première slide   --------------------------------------------

bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
# Ajout d'un titre
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Commentaires par segments'
# title_shape.text.font.underline = MSO_UNDERLINE.SINGLE_LINE
# title_shape.text.font.size = Pt(36)
# Ajout d'un cadre de texte auquel on ajoutera plusieurs paragraphes
tf = body_shape.text_frame

p = tf.add_paragraph()
p.text = 'Top 3 des commentaires de promoteurs : '
p.font.underline = MSO_UNDERLINE.SINGLE_LINE
p.level = 1
# On ajoute les 3 premiers commentaires de promoteurs
for i, x in enumerate(liste_comment_prom[0:4]):
    p = tf.add_paragraph()
    p.text = x
    p.font.bold = True
    p.level = 2
# On ajoute 1 commentaire d'un utilisateur passif
p = tf.add_paragraph()
p.text = 'Commentaire d\'un utilisateur passif : '
p.font.underline = MSO_UNDERLINE.SINGLE_LINE
p.level = 1
p = tf.add_paragraph()
p.text = liste_comment_pass[0]
p.font.bold = True
p.level = 2
# On ajoute les 1 commentaire d'un détracteur
p = tf.add_paragraph()
p.text = 'Commentaire de détracteur : '
p.font.underline = MSO_UNDERLINE.SINGLE_LINE
p.level = 1
p = tf.add_paragraph()
p.text = liste_comment_detr[0]
p.font.bold = True
p.level = 2

#--------------------    Deuxième slide qui ajoute le graphique  ---------------

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
#shapes = slide.shapes
img_path = 'data_graphiques/donut_chart.png'
# Taille de l'image
left = top = Inches(0.81)
# left = Inches(6)
height = Inches(7)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

#----------------------------   Troisième slide  -------------------------------

# troisième slide avec les 9 autres commentaires des promotteurs
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

tf = body_shape.text_frame
title_shape.text = 'Commentaires de promoteurs annexes : '
# title_shape.font.underline = MSO_UNDERLINE.SINGLE_LINE
# title_shape.font.size = Pt(36)
# On ajoute les 9 autres commentaires de promoteurs
for i, x in enumerate(liste_comment_prom[3:]):
    p = tf.add_paragraph()
    str_index = i+4
    p.text = 'Commentaire {} : '.format(str(str_index))
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.font.size = Pt(15)
    p.level = 1
    p = tf.add_paragraph()
    p.text = x
    p.font.size = Pt(15)
    p.font.bold = True
    p.level = 2

prs.save('test.pptx')